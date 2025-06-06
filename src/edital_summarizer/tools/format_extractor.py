import os
from typing import Type
from pydantic import BaseModel, Field
from crewai.tools import BaseTool
from ..utils.logger import get_logger
from pptx import Presentation
from openpyxl import load_workbook

logger = get_logger(__name__)

class FormatExtractorInput(BaseModel):
    """Input schema for FormatSpecificExtractor."""
    file_path: str = Field(..., description="Path to the file to be processed.")
    max_chars: int = Field(20000, description="Maximum number of characters to return.")

class FormatSpecificExtractor(BaseTool):
    name: str = "Format Specific Extractor"
    description: str = (
        "Extracts text from various file formats including PowerPoint and Excel. "
        "Handles format-specific processing and maintains document structure."
    )
    args_schema: Type[BaseModel] = FormatExtractorInput

    def _run(self, file_path: str, max_chars: int = 20000) -> str:
        """Extract text from various file formats."""
        if not os.path.exists(file_path):
            return f"Error: File not found at {file_path}"

        _, file_extension = os.path.splitext(file_path)
        file_extension = file_extension.lower()

        try:
            if file_extension in ['.ppt', '.pptx']:
                text = self._extract_from_ppt(file_path, max_chars)
            elif file_extension in ['.xls', '.xlsx']:
                text = self._extract_from_excel(file_path, max_chars)
            else:
                return f"Error: Unsupported file format {file_extension}"

            # Limit text size
            if len(text) > max_chars:
                text = text[:max_chars] + f"\n\n[Texto truncado em {max_chars} caracteres]"

            return text

        except Exception as e:
            logger.error(f"Error extracting text from {file_path}: {str(e)}")
            return f"Error: {str(e)}"

    def _extract_from_ppt(self, file_path: str, max_chars: int) -> str:
        """
        Extrai texto de um arquivo PPT/PPTX.
        
        Args:
            file_path: Caminho do arquivo
            max_chars: Número máximo de caracteres
            
        Returns:
            Texto extraído
        """
        try:
            prs = Presentation(file_path)
            text_parts = []
            
            # Extrai texto dos slides
            for slide in prs.slides:
                slide_text = []
                
                # Extrai texto dos shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                
                # Adiciona o texto do slide
                if slide_text:
                    text_parts.append("\n".join(slide_text))
            
            # Combina todo o texto
            full_text = "\n\n=== Novo Slide ===\n\n".join(text_parts)
            
            # Limita o tamanho
            if len(full_text) > max_chars:
                full_text = full_text[:max_chars] + "\n\n[Texto truncado]"
            
            return full_text
            
        except Exception as e:
            logger.error(f"Erro ao extrair texto do PPT {file_path}: {str(e)}")
            raise

    def _extract_from_excel(self, file_path: str, max_chars: int) -> str:
        """
        Extrai texto de um arquivo XLS/XLSX.
        
        Args:
            file_path: Caminho do arquivo
            max_chars: Número máximo de caracteres
            
        Returns:
            Texto extraído
        """
        try:
            wb = load_workbook(file_path, read_only=True, data_only=True)
            text_parts = []
            
            # Extrai texto das planilhas
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                sheet_text = []
                
                # Extrai texto das células
                for row in ws.rows:
                    row_text = []
                    for cell in row:
                        if cell.value:
                            row_text.append(str(cell.value))
                    if row_text:
                        sheet_text.append(" | ".join(row_text))
                
                # Adiciona o texto da planilha
                if sheet_text:
                    text_parts.append(f"\n=== Planilha: {sheet} ===\n")
                    text_parts.append("\n".join(sheet_text))
            
            # Combina todo o texto
            full_text = "\n\n".join(text_parts)
            
            # Limita o tamanho
            if len(full_text) > max_chars:
                full_text = full_text[:max_chars] + "\n\n[Texto truncado]"
            
            return full_text
            
        except Exception as e:
            logger.error(f"Erro ao extrair texto do Excel {file_path}: {str(e)}")
            raise 